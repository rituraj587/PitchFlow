import os
import io
import json
import httpx
import fitz  # PyMuPDF
from fastapi import FastAPI, Depends, HTTPException, status, UploadFile, File, Form, Header
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from sqlalchemy.orm import Session
from typing import Optional

# Import database, models, schemas, and auth
from backend.database import engine, SessionLocal, get_db
from backend.models import Base, User, GlobalConfig
from backend.schemas import (
    UserCreate, UserLogin, Token, UserResponse,
    AdminConfigUpdate, GlobalConfigResponse,
    PresentationSchema, Slide, SlideElement
)
from backend.auth import (
    get_password_hash, verify_password, create_access_token,
    get_current_user, get_current_admin
)

# python-pptx imports
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

app = FastAPI(title="Pitchflow API", version="1.0.0")

# CORS setup
origins = [
    "http://localhost:5173",
    "http://localhost:3000",
    "https://pitchflow.riturajshukla.tech",
]

app.add_middleware(
    CORSMiddleware,
    allow_origins=origins,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Startup DB Seeding
@app.on_event("startup")
def startup_event():
    # Create database tables
    Base.metadata.create_all(bind=engine)
    
    db = SessionLocal()
    try:
        # 1. Seed GlobalConfig
        config = db.query(GlobalConfig).filter(GlobalConfig.id == 1).first()
        if not config:
            new_config = GlobalConfig(
                id=1,
                max_pdf_pages=10,
                max_slides_allowed=12,
                allow_signups=True
            )
            db.add(new_config)
            db.commit()

        # 2. Seed Default Admin User
        admin_email = os.getenv("ADMIN_EMAIL", "admin@pitchflow.com")
        admin_password = os.getenv("ADMIN_PASSWORD", "admin123")
        admin = db.query(User).filter(User.email == admin_email).first()
        if not admin:
            new_admin = User(
                email=admin_email,
                hashed_password=get_password_hash(admin_password),
                is_admin=True
            )
            db.add(new_admin)
            db.commit()
            print(f"Default admin seeded: {admin_email} / {admin_password}")

        # Seed Admin Rupali
        rupali_email = "rupali@pitchflow.com"
        rupali = db.query(User).filter(User.email == rupali_email).first()
        if not rupali:
            new_rupali = User(
                email=rupali_email,
                hashed_password=get_password_hash("rupali123"),
                is_admin=True
            )
            db.add(new_rupali)
            db.commit()
            print(f"Admin Rupali seeded: {rupali_email} / rupali123")

        # 3. Seed Default Standard User
        user_email = os.getenv("USER_EMAIL", "user@pitchflow.com")
        user_password = os.getenv("USER_PASSWORD", "user123")
        user = db.query(User).filter(User.email == user_email).first()
        if not user:
            new_user = User(
                email=user_email,
                hashed_password=get_password_hash(user_password),
                is_admin=False
            )
            db.add(new_user)
            db.commit()
            print(f"Default standard user seeded: {user_email} / {user_password}")
    finally:
        db.close()


# --- Auth Routes ---

@app.post("/api/auth/signup", response_model=UserResponse, status_code=status.HTTP_201_CREATED)
def signup(user_in: UserCreate, db: Session = Depends(get_db)):
    config = db.query(GlobalConfig).filter(GlobalConfig.id == 1).first()
    if not config or not config.allow_signups:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="Signups are currently disabled by administrator"
        )
    
    existing_user = db.query(User).filter(User.email == user_in.email).first()
    if existing_user:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Email already registered"
        )
    
    new_user = User(
        email=user_in.email,
        hashed_password=get_password_hash(user_in.password),
        is_admin=False
    )
    db.add(new_user)
    db.commit()
    db.refresh(new_user)
    return new_user


@app.post("/api/auth/login", response_model=Token)
def login(credentials: UserLogin, db: Session = Depends(get_db)):
    user = db.query(User).filter(User.email == credentials.email).first()
    if not user or not verify_password(credentials.password, user.hashed_password):
        raise HTTPException(
            status_code=status.HTTP_401_UNAUTHORIZED,
            detail="Invalid email or password",
            headers={"WWW-Authenticate": "Bearer"},
        )
    
    # Create access token and include email and is_admin
    access_token = create_access_token(data={"sub": user.email, "is_admin": user.is_admin})
    return {"access_token": access_token, "token_type": "bearer"}


# --- Admin Routes ---

@app.get("/api/admin/config", response_model=GlobalConfigResponse)
def get_config(db: Session = Depends(get_db), current_admin: User = Depends(get_current_admin)):
    config = db.query(GlobalConfig).filter(GlobalConfig.id == 1).first()
    if not config:
        raise HTTPException(status_code=404, detail="Configuration not found")
    return config


@app.post("/api/admin/config", response_model=GlobalConfigResponse)
def update_config(config_in: AdminConfigUpdate, db: Session = Depends(get_db), current_admin: User = Depends(get_current_admin)):
    config = db.query(GlobalConfig).filter(GlobalConfig.id == 1).first()
    if not config:
        raise HTTPException(status_code=404, detail="Configuration not found")
    
    if config_in.max_pdf_pages is not None:
        config.max_pdf_pages = config_in.max_pdf_pages
    if config_in.max_slides_allowed is not None:
        config.max_slides_allowed = config_in.max_slides_allowed
    if config_in.allow_signups is not None:
        config.allow_signups = config_in.allow_signups
        
    db.commit()
    db.refresh(config)
    return config


# --- LLM API Call Helpers ---

async def call_openai(api_key: str, pdf_text: str, max_slides: int) -> dict:
    url = "https://api.openai.com/v1/chat/completions"
    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json"
    }
    system_prompt = (
        "You are an expert presentation designer. Analyze the provided PDF text and summarize it into a professionally structured, detailed presentation in JSON format.\n"
        f"Generate at most {max_slides} slides.\n"
        "To make the presentation highly informative and elite, write comprehensive, explanatory bullet points. Avoid single words or short phrases. Each slide must contain between 4 and 6 descriptive elements (e.g. 1 subtitle/overview and 3-5 descriptive bullets) that thoroughly explain the details, facts, definitions, and logic from the source document. Each bullet point sentence MUST be informative but strictly limited to at most 120 characters.\n"
        "The output MUST be a JSON object matching this schema exactly:\n"
        "{\n"
        '  "title": "Overall Presentation Title",\n'
        '  "slides": [\n'
        "    {\n"
        '      "title": "Slide Title",\n'
        '      "elements": [\n'
        '        {"type": "subtitle", "content": "Detailed subtitle or high-level outline summary (max 120 chars)"},\n'
        '        {"type": "bullet_point", "content": "Explanatory point 1 (max 120 chars)"},\n'
        '        {"type": "bullet_point", "content": "Explanatory point 2 (max 120 chars)"}\n'
        "      ]\n"
        "    }\n"
        "  ]\n"
        "}\n"
        "Do not wrap your response in markdown code blocks. Output raw JSON only."
    )
    user_prompt = f"PDF Content:\n{pdf_text}"
    
    payload = {
        "model": "gpt-4o-mini",
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ],
        "response_format": {"type": "json_object"},
        "temperature": 0.3
    }
    
    async with httpx.AsyncClient(timeout=90.0) as client:
        response = await client.post(url, headers=headers, json=payload)
        if response.status_code != 200:
            raise HTTPException(
                status_code=response.status_code,
                detail=f"OpenAI API Error: {response.text}"
            )
        data = response.json()
        content = data["choices"][0]["message"]["content"]
        return json.loads(content)


async def call_gemini(api_key: str, pdf_text: str, max_slides: int) -> dict:
    url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-1.5-flash:generateContent?key={api_key}"
    headers = {
        "Content-Type": "application/json"
    }
    system_prompt = (
        "You are an expert presentation designer. Analyze the provided PDF text and summarize it into a professionally structured, detailed presentation in JSON format.\n"
        f"Generate at most {max_slides} slides.\n"
        "To make the presentation highly informative and elite, write comprehensive, explanatory bullet points. Avoid single words or short phrases. Each slide must contain between 4 and 6 descriptive elements (e.g. 1 subtitle/overview and 3-5 descriptive bullets) that thoroughly explain the details, facts, definitions, and logic from the source document. Each bullet point sentence MUST be informative but strictly limited to at most 120 characters.\n"
        "The output MUST be a JSON object matching this schema exactly:\n"
        "{\n"
        '  "title": "Overall Presentation Title",\n'
        '  "slides": [\n'
        "    {\n"
        '      "title": "Slide Title",\n'
        '      "elements": [\n'
        '        {"type": "subtitle", "content": "Detailed subtitle or high-level outline summary (max 120 chars)"},\n'
        '        {"type": "bullet_point", "content": "Explanatory point 1 (max 120 chars)"},\n'
        '        {"type": "bullet_point", "content": "Explanatory point 2 (max 120 chars)"}\n'
        "      ]\n"
        "    }\n"
        "  ]\n"
        "}\n"
        "Do not wrap your response in markdown code blocks. Output raw JSON only."
    )
    
    payload = {
        "contents": [
            {
                "parts": [
                    {"text": f"PDF Content:\n{pdf_text}"}
                ]
            }
        ],
        "systemInstruction": {
            "parts": [
                {"text": system_prompt}
            ]
        },
        "generationConfig": {
            "responseMimeType": "application/json",
            "temperature": 0.3
        }
    }
    
    async with httpx.AsyncClient(timeout=90.0) as client:
        response = await client.post(url, headers=headers, json=payload)
        if response.status_code != 200:
            raise HTTPException(
                status_code=response.status_code,
                detail=f"Gemini API Error: {response.text}"
            )
        data = response.json()
        content = data["candidates"][0]["content"]["parts"][0]["text"]
        return json.loads(content)


# --- PPTX Generation Engine ---

def create_presentation_file(data: PresentationSchema) -> io.BytesIO:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    
    # Premium Modern Theme Colors (Corporate Tech Standard - Dark Slate & Cobalt Blue)
    bg_color = RGBColor(11, 15, 25)        # Slate 950 (#0B0F19)
    accent_color = RGBColor(37, 99, 235)    # Cobalt Blue (#2563EB)
    text_light = RGBColor(248, 250, 252)    # Slate 50 (#F8FAFC)
    text_muted = RGBColor(148, 163, 184)    # Slate 400 (#94A3B8)
    
    # 1. Title Slide
    slide = prs.slides.add_slide(blank_layout)
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color
    
    # Left decorative bar
    left_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), Inches(1.8), Inches(0.15), Inches(3.9)
    )
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = accent_color
    left_bar.line.color.rgb = accent_color
    
    # Title & Subtitle textbox
    title_box = slide.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(11.0), Inches(3.9))
    tf = title_box.text_frame
    tf.word_wrap = True
    # Zero-Margin Containment
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    
    p_title = tf.paragraphs[0]
    p_title.text = data.title
    p_title.font.name = "Segoe UI"
    p_title.font.size = Pt(44)
    p_title.font.bold = True
    p_title.font.color.rgb = text_light
    p_title.space_after = Pt(16)
    
    p_sub = tf.add_paragraph()
    p_sub.text = "Generated by Pitchflow • Professional Document Summarization Engine"
    p_sub.font.name = "Segoe UI"
    p_sub.font.size = Pt(16)
    p_sub.font.color.rgb = text_muted
    
    # 2. Content Slides
    for slide_data in data.slides:
        slide = prs.slides.add_slide(blank_layout)
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
        
        # Widescreen top color accent bar
        top_strip = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(13.333), Inches(0.1)
        )
        top_strip.fill.solid()
        top_strip.fill.fore_color.rgb = accent_color
        top_strip.line.color.rgb = accent_color
        
        # Slide Title
        header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.5), Inches(1.0))
        tf_header = header_box.text_frame
        tf_header.word_wrap = True
        tf_header.margin_left = Inches(0)
        tf_header.margin_right = Inches(0)
        tf_header.margin_top = Inches(0)
        tf_header.margin_bottom = Inches(0)
        p_hdr = tf_header.paragraphs[0]
        p_hdr.text = slide_data.title
        p_hdr.font.name = "Segoe UI"
        p_hdr.font.size = Pt(30)
        p_hdr.font.bold = True
        p_hdr.font.color.rgb = text_light
        
        # Split Elements: subtitle and content points
        content_elements = [e for e in slide_data.elements if e.type.lower() != "subtitle"]
        subtitle_element = next((e for e in slide_data.elements if e.type.lower() == "subtitle"), None)
        
        # Draw subtitle if present
        if subtitle_element:
            sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.5))
            tf_sub = sub_box.text_frame
            tf_sub.word_wrap = True
            tf_sub.margin_left = Inches(0)
            tf_sub.margin_right = Inches(0)
            tf_sub.margin_top = Inches(0)
            tf_sub.margin_bottom = Inches(0)
            p_sub = tf_sub.paragraphs[0]
            p_sub.text = subtitle_element.content[:120] # String Slicing
            p_sub.font.name = "Segoe UI"
            p_sub.font.size = Pt(16)
            p_sub.font.italic = True
            p_sub.font.color.rgb = text_muted
            
        # Draw columns based on element count
        top_y = Inches(1.8)
        height_y = Inches(5.0)
        
        if len(content_elements) > 2:
            # Defensive Grid Segmentation: Left col (X=1.0, W=4.5), Right col (X=6.5, W=5.8)
            left_list = content_elements[:3]      # Double Slicing (Max 3)
            right_list = content_elements[3:6]    # Double Slicing (Max 3)
            
            # Left column box
            left_box = slide.shapes.add_textbox(Inches(1.0), top_y, Inches(4.5), height_y)
            tf_left = left_box.text_frame
            tf_left.word_wrap = True
            tf_left.margin_left = Inches(0)
            tf_left.margin_right = Inches(0)
            tf_left.margin_top = Inches(0)
            tf_left.margin_bottom = Inches(0)
            
            for idx, elem in enumerate(left_list):
                p = tf_left.paragraphs[0] if idx == 0 else tf_left.add_paragraph()
                p.text = elem.content[:120]  # String Slicing (Max 120 chars)
                p.font.name = "Segoe UI"
                p.font.size = Pt(15)
                p.font.color.rgb = text_light
                p.level = 0
                p.space_after = Pt(12)
                
            # Right column box
            if right_list:
                right_box = slide.shapes.add_textbox(Inches(6.5), top_y, Inches(5.8), height_y)
                tf_right = right_box.text_frame
                tf_right.word_wrap = True
                tf_right.margin_left = Inches(0)
                tf_right.margin_right = Inches(0)
                tf_right.margin_top = Inches(0)
                tf_right.margin_bottom = Inches(0)
                
                for idx, elem in enumerate(right_list):
                    p = tf_right.paragraphs[0] if idx == 0 else tf_right.add_paragraph()
                    p.text = elem.content[:120]  # String Slicing (Max 120 chars)
                    p.font.name = "Segoe UI"
                    p.font.size = Pt(15)
                    p.font.color.rgb = text_light
                    p.level = 0
                    p.space_after = Pt(12)
        else:
            # Single column box
            box = slide.shapes.add_textbox(Inches(1.0), top_y, Inches(11.0), height_y)
            tf_box = box.text_frame
            tf_box.word_wrap = True
            tf_box.margin_left = Inches(0)
            tf_box.margin_right = Inches(0)
            tf_box.margin_top = Inches(0)
            tf_box.margin_bottom = Inches(0)
            
            for idx, elem in enumerate(content_elements[:3]):  # Double Slicing (Max 3)
                p = tf_box.paragraphs[0] if idx == 0 else tf_box.add_paragraph()
                p.text = elem.content[:120]  # String Slicing (Max 120 chars)
                p.font.name = "Segoe UI"
                p.font.size = Pt(16)
                p.font.color.rgb = text_light
                p.level = 0
                p.space_after = Pt(14)
                
    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io


# --- Business Engine Route ---

@app.post("/api/generate-presentation")
async def generate_presentation(
    file: UploadFile = File(...),
    model_choice: str = Form(...),
    x_openai_key: Optional[str] = Header(None, alias="X-OpenAI-Key"),
    x_gemini_key: Optional[str] = Header(None, alias="X-Gemini-Key"),
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    # Fetch global rules
    config = db.query(GlobalConfig).filter(GlobalConfig.id == 1).first()
    if not config:
        raise HTTPException(status_code=500, detail="Global application configuration missing")
    
    # 1. Read PDF & Extract Text
    try:
        pdf_bytes = await file.read()
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        page_count = len(doc)
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Failed to parse PDF file: {str(e)}")
    
    # Validate PDF pages limits
    if page_count > config.max_pdf_pages:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Uploaded PDF exceeds maximum page limit of {config.max_pdf_pages} (file has {page_count} pages)."
        )
    
    # Extract text from all pages
    extracted_text = ""
    for idx in range(page_count):
        extracted_text += f"\n--- Page {idx + 1} ---\n"
        extracted_text += doc[idx].get_text()
        
    doc.close()
    
    # Check for empty content
    if not extracted_text.strip():
        raise HTTPException(status_code=400, detail="Uploaded PDF contains no extractable text.")
        
    # 2. Invoke LLM dynamically
    raw_response = None
    if model_choice == "openai":
        if not x_openai_key or not x_openai_key.strip():
            raise HTTPException(status_code=400, detail="OpenAI API key missing in X-OpenAI-Key header")
        raw_response = await call_openai(x_openai_key.strip(), extracted_text, config.max_slides_allowed)
    elif model_choice == "gemini":
        if not x_gemini_key or not x_gemini_key.strip():
            raise HTTPException(status_code=400, detail="Gemini API key missing in X-Gemini-Key header")
        raw_response = await call_gemini(x_gemini_key.strip(), extracted_text, config.max_slides_allowed)
    else:
        raise HTTPException(status_code=400, detail="Invalid model_choice. Select 'openai' or 'gemini'.")
    
    # 3. Parse and Validate Presentation Structure
    try:
        presentation_data = PresentationSchema(**raw_response)
    except Exception as e:
        raise HTTPException(
            status_code=500,
            detail=f"LLM failed to output correct structure. Details: {str(e)}"
        )
        
    # Slice slide count to strictly observe limits
    presentation_data.slides = presentation_data.slides[:config.max_slides_allowed]
    
    # 4. Generate PPTX
    try:
        pptx_buffer = create_presentation_file(presentation_data)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to generate PPTX file: {str(e)}")
        
    # Return streaming download
    safe_filename = file.filename.rsplit(".", 1)[0] + "_pitchflow.pptx"
    return StreamingResponse(
        pptx_buffer,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{safe_filename}"'}
    )
